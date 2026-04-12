import io
import logging

import requests
from copy import deepcopy
from lxml import etree
from PIL import Image
from pptx import Presentation
from pptx.oxml.ns import qn

logger = logging.getLogger("slide-generator")

# Template slide indices (0-based)
TITLE_IDX = 0
TEXT_ONLY_IDX = 1
HERO_IMAGE_IDX = 2
TWO_COL_IDX = 3
THREE_COL_IDX = 4

LAYOUT_TO_TEMPLATE = {
    "text_only": TEXT_ONLY_IDX,
    "hero_image": HERO_IMAGE_IDX,
    "two_column": TWO_COL_IDX,
    "three_column": THREE_COL_IDX,
}


def generate_pptx(template_path, sections, output_path):
    """Generate a PPTX by cloning template slides and injecting content."""
    prs = Presentation(template_path)
    num_template_slides = len(prs.slides)

    # 2. Content slides
    for section in sections:
        for slide_data in section.get("slides", []):
            layout = slide_data.get("layout", "text_only")
            tmpl_idx = LAYOUT_TO_TEMPLATE.get(layout, TEXT_ONLY_IDX)
            new_slide = _clone_slide(prs, tmpl_idx)
            _fill_content_slide(new_slide, slide_data, layout)

    # 3. Remove original template slides (always delete index 0)
    for _ in range(num_template_slides):
        _delete_slide(prs, 0)

    prs.save(output_path)


# ---------------------------------------------------------------------------
#  Slide cloning
# ---------------------------------------------------------------------------

def _clone_slide(prs, source_idx):
    """Clone slide at source_idx; new slide is appended at the end."""
    source = prs.slides[source_idx]
    new_slide = prs.slides.add_slide(source.slide_layout)

    # Map old rIds → new rIds for image relationships
    rId_map = {}
    for key, rel in source.part.rels.items():
        if "image" in rel.reltype:
            new_rId = new_slide.part.relate_to(rel.target_part, rel.reltype)
            rId_map[key] = new_rId

    # Replace shape tree contents
    new_spTree = new_slide.shapes._spTree
    _clear_shapes(new_spTree)

    for child in source.shapes._spTree:
        if _is_shape_element(child):
            new_elem = deepcopy(child)
            for blip in new_elem.iter(qn("a:blip")):
                old_embed = blip.get(qn("r:embed"))
                if old_embed and old_embed in rId_map:
                    blip.set(qn("r:embed"), rId_map[old_embed])
            new_spTree.append(new_elem)

    return new_slide


def _clear_shapes(spTree):
    keep = {"nvGrpSpPr", "grpSpPr"}
    for child in list(spTree):
        if etree.QName(child.tag).localname not in keep:
            spTree.remove(child)


def _is_shape_element(elem):
    return etree.QName(elem.tag).localname not in {"nvGrpSpPr", "grpSpPr"}


def _delete_slide(prs, index):
    sldIdLst = prs.slides._sldIdLst
    sldId = list(sldIdLst)[index]
    rId = sldId.get(qn("r:id"))
    prs.part.drop_rel(rId)
    sldIdLst.remove(sldId)


# ---------------------------------------------------------------------------
#  Content injection – Title slide
# ---------------------------------------------------------------------------

def _fill_title_slide(slide, title):
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name == "TextBox 41":
            _set_run_text(shape, title)
            break


# ---------------------------------------------------------------------------
#  Content injection – Content slides
# ---------------------------------------------------------------------------

def _fill_content_slide(slide, data, layout):
    title = data.get("title", "")
    key_points = data.get("key_points", [])
    images = data.get("images", [])
    transcript = data.get("transcript", "")

    _set_placeholder_title(slide, title)
    _set_key_points(slide, key_points)

    if layout in ("hero_image", "two_column", "three_column") and images:
        _replace_images(slide, images)
        _update_captions(slide, images)

    if transcript:
        _set_speaker_notes(slide, transcript)


def _set_placeholder_title(slide, title):
    for shape in slide.shapes:
        if shape.name == "Title 1" and shape.has_text_frame:
            _set_run_text(shape, title)
            return


def _set_key_points(slide, key_points):
    for shape in slide.shapes:
        if "Content Placeholder" not in shape.name:
            continue
        if not shape.has_text_frame:
            continue

        txBody = shape.text_frame._txBody
        existing = list(txBody.iterchildren(qn("a:p")))
        if not existing or not key_points:
            return

        template_p = deepcopy(existing[0])
        for p in existing:
            txBody.remove(p)

        for kp in key_points:
            new_p = deepcopy(template_p)
            for t in new_p.iter(qn("a:t")):
                t.text = kp
                break
            txBody.append(new_p)
        return


def _replace_images(slide, images):
    """Remove template images and insert downloaded ones at the same positions."""
    pic_shapes = sorted(
        [s for s in slide.shapes if s.shape_type == 13],
        key=lambda s: s.left,
    )

    for pic_shape, img_data in zip(pic_shapes, images):
        left, top = pic_shape.left, pic_shape.top
        width, height = pic_shape.width, pic_shape.height

        pic_shape._element.getparent().remove(pic_shape._element)

        img_stream = _download_and_crop(img_data["url"], width, height)
        if img_stream:
            slide.shapes.add_picture(img_stream, left, top, width, height)


def _update_captions(slide, images):
    """Set caption text in TextBox shapes located in the lower region of the slide."""
    captions = sorted(
        [
            s for s in slide.shapes
            if s.shape_type == 17 and s.has_text_frame and s.top > 5_000_000
        ],
        key=lambda s: s.left,
    )
    for cap_shape, img_data in zip(captions, images):
        _set_run_text(cap_shape, img_data.get("caption", ""))


def _set_speaker_notes(slide, text):
    notes = slide.notes_slide
    notes.notes_text_frame.text = text


# ---------------------------------------------------------------------------
#  Helpers
# ---------------------------------------------------------------------------

def _set_run_text(shape, text):
    """Set text of the first run in the first paragraph, preserving formatting."""
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            run.text = text
            return


def _download_and_crop(url, target_w_emu, target_h_emu):
    """Download an image, centre-crop to the template slot's aspect ratio."""
    try:
        resp = requests.get(url, timeout=30)
        resp.raise_for_status()
        img = Image.open(io.BytesIO(resp.content)).convert("RGB")

        target_ratio = target_w_emu / target_h_emu
        img_ratio = img.width / img.height

        if abs(img_ratio - target_ratio) > 0.01:
            if img_ratio > target_ratio:
                new_w = int(img.height * target_ratio)
                left = (img.width - new_w) // 2
                img = img.crop((left, 0, left + new_w, img.height))
            else:
                new_h = int(img.width / target_ratio)
                top = (img.height - new_h) // 2
                img = img.crop((0, top, img.width, top + new_h))

        buf = io.BytesIO()
        img.save(buf, format="PNG")
        buf.seek(0)
        return buf
    except Exception:
        logger.exception("Image download/crop failed for %s", url)
        return None
